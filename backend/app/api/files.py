"""ClarityLoop – File Upload API"""
from __future__ import annotations
import io
import json
from pathlib import Path
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File as FastFile
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from app.db.session import get_db
from app.models.file import UploadedFile
from app.core.config import get_settings
from app.core.security import validate_file_upload, sanitize_filename, unique_safe_filename, detect_prompt_injection
from app.schemas import UploadedFileOut

router = APIRouter()
settings = get_settings()
DEFAULT_USER = "default-user"
UPLOAD_DIR = Path("uploads")


def _parse_text(content: bytes, ext: str) -> tuple[str, dict]:
    """Parse file content based on extension. Returns (text, metadata)."""
    meta: dict = {}
    try:
        if ext in (".csv",):
            import pandas as pd
            df = pd.read_csv(io.BytesIO(content))
            meta = {
                "rows": len(df),
                "columns": list(df.columns.tolist()),
                "col_count": len(df.columns),
                "missing_values": df.isnull().sum().to_dict(),
                "dtypes": df.dtypes.astype(str).to_dict(),
            }
            return df.head(100).to_string(), meta
        elif ext in (".xlsx", ".xls"):
            import pandas as pd
            df = pd.read_excel(io.BytesIO(content))
            meta = {
                "rows": len(df),
                "columns": list(df.columns.tolist()),
                "col_count": len(df.columns),
            }
            return df.head(100).to_string(), meta
        elif ext == ".pdf":
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages[:20])
            meta = {"pages": len(reader.pages)}
            return text, meta
        elif ext in (".docx",):
            import docx
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs)
            return text, meta
        elif ext == ".pptx":
            from pptx import Presentation
            presentation = Presentation(io.BytesIO(content))
            slides = []
            for index, slide in enumerate(presentation.slides, start=1):
                slide_text = "\n".join(
                    shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
                )
                slides.append(f"Slide {index}:\n{slide_text}")
            meta = {"slides": len(presentation.slides)}
            return "\n\n".join(slides), meta
        elif ext in (".json",):
            data = json.loads(content.decode("utf-8", errors="replace"))
            return json.dumps(data, indent=2)[:5000], meta
        else:
            # Plain text, markdown, code
            return content.decode("utf-8", errors="replace"), meta
    except Exception as e:
        return "", {"parse_error": str(e)}


@router.post("/upload/{project_id}", response_model=UploadedFileOut, status_code=201)
async def upload_file(
    project_id: str,
    file: UploadFile = FastFile(...),
    db: AsyncSession = Depends(get_db),
):
    content = await file.read()
    ext = Path(file.filename or "").suffix.lower()

    # Validate
    ok, err = validate_file_upload(
        filename=file.filename or "upload",
        mime_type=file.content_type or "application/octet-stream",
        size_bytes=len(content),
        max_bytes=settings.max_upload_bytes,
    )
    if not ok:
        raise HTTPException(400, err)

    # Sanitize filename
    safe_name = unique_safe_filename(file.filename or "upload", UPLOAD_DIR)

    # Save to disk
    UPLOAD_DIR.mkdir(exist_ok=True)
    file_path = UPLOAD_DIR / safe_name
    file_path.write_bytes(content)

    # Parse content (for text files)
    is_image = ext in (".png", ".jpg", ".jpeg", ".webp")
    parsed_text = ""
    meta: dict = {}
    parse_error = None

    if not is_image:
        parsed_text, meta = _parse_text(content, ext)
        if "parse_error" in meta:
            parse_error = meta.pop("parse_error")

        # Prompt injection detection
        if parsed_text:
            detected, matched = detect_prompt_injection(parsed_text)
            if detected:
                parsed_text = f"[Content sanitized: potential injection pattern detected near '{matched}']\n" + parsed_text

    record = UploadedFile(
        project_id=project_id,
        user_id=DEFAULT_USER,
        original_filename=file.filename or "upload",
        safe_filename=safe_name,
        storage_path=str(file_path),
        mime_type=file.content_type or "application/octet-stream",
        file_extension=ext,
        size_bytes=len(content),
        parsed_content=parsed_text[:50000] if parsed_text else None,
        parse_error=parse_error,
        is_image=is_image,
        row_count=meta.get("rows"),
        col_count=meta.get("col_count"),
        schema_json=json.dumps(meta) if meta else None,
    )
    db.add(record)
    await db.flush()
    await db.refresh(record)
    return record

@router.get("/{project_id}", response_model=list[UploadedFileOut])
async def list_files(project_id: str, db: AsyncSession = Depends(get_db)):
    result = await db.execute(select(UploadedFile).where(UploadedFile.project_id == project_id))
    return result.scalars().all()

@router.get("/detail/{file_id}", response_model=UploadedFileOut)
async def get_file(file_id: str, db: AsyncSession = Depends(get_db)):
    result = await db.execute(select(UploadedFile).where(UploadedFile.id == file_id))
    file_rec = result.scalar_one_or_none()
    if not file_rec:
        raise HTTPException(404, "File not found")
    return file_rec
